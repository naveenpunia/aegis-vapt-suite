from __future__ import annotations
import csv
import json
from pathlib import Path
from typing import Any
from jinja2 import Environment, FileSystemLoader, select_autoescape
from tabulate import tabulate
from utils.helpers import calculate_risk_score, color_status, ensure_output_dir
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
except ModuleNotFoundError:
    colors = None
    A4 = None
    getSampleStyleSheet = None
    Paragraph = None
    SimpleDocTemplate = None
    Spacer = None
    Table = None
    TableStyle = None

class ReportGenerator:

    def __init__(self, output_dir: str='output', template_dir: str='templates') -> None:
        self.output_dir = ensure_output_dir(output_dir)
        self.template_dir = Path(template_dir)

    def enrich_results(self, results: dict[str, Any]) -> dict[str, Any]:
        risk_score, risk_severity, recommendations = calculate_risk_score(results)
        enriched = dict(results)
        enriched['risk_score'] = risk_score
        enriched['risk_severity'] = risk_severity
        enriched['recommendations'] = recommendations
        return enriched

    def print_console(self, results: dict[str, Any]) -> None:
        enriched = self.enrich_results(results)
        print('\n=== MINI PENETRATION TESTING TOOLKIT REPORT ===')
        print(f"Target URL: {enriched['target_url']}")
        print(f"Scan Time : {enriched['scan_time']}")
        print(f"Risk Score: {enriched['risk_score']}/100 ({color_status(enriched['risk_severity'], enriched['risk_severity'])})")
        url_scan = enriched.get('url_scan', {})
        print('\n[URL Scanner]')
        print(f"Valid URL: {url_scan.get('valid')}")
        print(f"HTTP Status: {url_scan.get('status_code')} {url_scan.get('reason')}")
        print(f"Final URL: {url_scan.get('final_url')}")
        print(f"Server: {url_scan.get('server')}")
        print(f"Missing Security Headers: {', '.join(url_scan.get('missing_security_headers', [])) or 'None'}")
        print('\n[Port Scanner]')
        print(tabulate(enriched.get('port_scan', []), headers='keys', tablefmt='grid'))
        print('\n[SQL Injection Tester]')
        print(tabulate(enriched.get('sql_injection', []), headers='keys', tablefmt='grid'))
        print('\n[XSS Detection]')
        print(tabulate(enriched.get('xss', []), headers='keys', tablefmt='grid'))
        print('\n[SSL Checker]')
        ssl_status = enriched.get('ssl', {})
        for key, value in ssl_status.items():
            print(f'{key}: {value}')
        print('\n[Recommendations]')
        for item in enriched.get('recommendations', []):
            print(f'- {item}')

    def save_all(self, results: dict[str, Any], base_name: str) -> dict[str, str]:
        enriched = self.enrich_results(results)
        paths = {'txt': str(self.save_txt(enriched, base_name)), 'html': str(self.save_html(enriched, base_name)), 'csv': str(self.save_csv(enriched, base_name))}
        ppt_path = self.save_ppt(enriched, base_name)
        if ppt_path:
            paths['ppt'] = str(ppt_path)
        pdf_path = self.save_pdf(enriched, base_name)
        if pdf_path:
            paths['pdf'] = str(pdf_path)
        return paths

    def save_txt(self, results: dict[str, Any], base_name: str) -> Path:
        path = self.output_dir / f'{base_name}.txt'
        recon_data = results.get('recon_scan', {})
        subdomains_txt = "\n".join([f"  - {item['subdomain']} ({item['ip']})" for item in recon_data.get('subdomains', [])]) or "  None"
        techs_txt = ", ".join(recon_data.get('technologies', [])) or "None"
        
        lines = [
            'ADVANCED VAPT SUITE ASSESSMENT REPORT', 
            f"Target URL: {results['target_url']}", 
            f"Scan Time: {results['scan_time']}", 
            f"Risk Score: {results['risk_score']}/100 ({results['risk_severity']})", 
            '', 
            'URL Scan:', json.dumps(results.get('url_scan', {}), indent=2), 
            '',
            'Reconnaissance (Technologies & Subdomains):',
            f"  Technologies: {techs_txt}",
            "  Discovered Subdomains:", subdomains_txt,
            '',
            'Port Scan:', tabulate(results.get('port_scan', []), headers='keys'), 
            '', 
            'SQL Injection:', tabulate(results.get('sql_injection', []), headers='keys'), 
            '', 
            'XSS:', tabulate(results.get('xss', []), headers='keys'), 
            '', 
            'SSL:', json.dumps(results.get('ssl', {}), indent=2), 
            '', 
            'Recommendations:', '\r\n'.join((f'- {item}' for item in results.get('recommendations', [])))
        ]
        path.write_text('\r\n'.join(lines), encoding='utf-8-sig')
        return path

    def save_html(self, results: dict[str, Any], base_name: str) -> Path:
        path = self.output_dir / f'{base_name}.html'
        environment = Environment(loader=FileSystemLoader(self.template_dir), autoescape=select_autoescape(['html', 'xml']))
        template = environment.get_template('report.html')
        path.write_text(template.render(**results), encoding='utf-8')
        return path

    def save_csv(self, results: dict[str, Any], base_name: str) -> Path:
        path = self.output_dir / f'{base_name}.csv'
        with path.open('w', newline='', encoding='utf-8') as file_obj:
            writer = csv.writer(file_obj)
            writer.writerow(['section', 'field_1', 'field_2', 'field_3', 'field_4'])
            recon_data = results.get('recon_scan', {})
            for tech in recon_data.get('technologies', []):
                writer.writerow(['recon_tech', tech, '', '', ''])
            for item in recon_data.get('subdomains', []):
                writer.writerow(['recon_subdomain', item['subdomain'], item['ip'], '', ''])
            for row in results.get('port_scan', []):
                writer.writerow(['port_scan', row['port'], row['service'], row['status'], ''])
            for row in results.get('sql_injection', []):
                writer.writerow(['sql_injection', row['parameter'], row['payload'], row['vulnerable'], row['evidence']])
            for row in results.get('xss', []):
                writer.writerow(['xss', row['parameter'], row['payload'], row['vulnerable'], row['evidence']])
            ssl_status = results.get('ssl', {})
            writer.writerow(['ssl', ssl_status.get('host'), ssl_status.get('tls_version'), ssl_status.get('valid'), ssl_status.get('expiry_date')])
        return path

    def save_ppt(self, results: dict[str, Any], base_name: str) -> Path | None:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ModuleNotFoundError:
            note_path = self.output_dir / f'{base_name}_ppt_not_created.txt'
            note_path.write_text('PPTX report was not created because python-pptx is not installed. Run: pip install -r requirements.txt', encoding='utf-8-sig')
            return None
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        dark_bg = RGBColor(9, 13, 22)
        light_bg = RGBColor(248, 250, 252)
        text_light = RGBColor(241, 245, 249)
        text_dark = RGBColor(15, 23, 42)
        accent_blue = RGBColor(14, 165, 233)
        accent_grey = RGBColor(148, 163, 184)
        def set_bg(slide, color):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color
        def add_slide_header(slide, title_text, dark_theme=False):
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = 'Arial'
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = text_light if dark_theme else text_dark
        slide_layout = prs.slide_layouts[6]
        slide1 = prs.slides.add_slide(slide_layout)
        set_bg(slide1, dark_bg)
        title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = "SECURITY ASSESSMENT REPORT"
        p.font.name = 'Arial'
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = accent_blue
        p2 = tf.add_paragraph()
        p2.text = "Mini Penetration Testing Toolkit Scan Results"
        p2.font.name = 'Arial'
        p2.font.size = Pt(22)
        p2.font.color.rgb = text_light
        p2.space_before = Pt(14)
        p3 = tf.add_paragraph()
        p3.text = f"Target URL: {results['target_url']}\nScan Time: {results['scan_time']}"
        p3.font.name = 'Arial'
        p3.font.size = Pt(14)
        p3.font.color.rgb = accent_grey
        p3.space_before = Pt(36)
        slide2 = prs.slides.add_slide(slide_layout)
        set_bg(slide2, light_bg)
        add_slide_header(slide2, "Executive Summary")
        summary_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.833), Inches(5.0))
        tf2 = summary_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p = tf2.paragraphs[0]
        p.text = "Overall Security Assessment Summary"
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = text_dark
        p2 = tf2.add_paragraph()
        p2.text = f"- Risk Score: {results['risk_score']} / 100"
        p2.font.name = 'Arial'
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = text_dark
        p2.space_before = Pt(12)
        p3 = tf2.add_paragraph()
        p3.text = f"- Risk Severity: {results['risk_severity']}"
        p3.font.name = 'Arial'
        p3.font.size = Pt(16)
        p3.font.bold = True
        p3.font.color.rgb = text_dark
        p3.space_before = Pt(8)
        p4 = tf2.add_paragraph()
        p4.text = "- Summary of Scans Performed:\n" \
                  f"  - URL Health & Headers Check: Completed\n" \
                  f"  - Open Port Probe: Checked {len(results.get('port_scan', []))} ports\n" \
                  f"  - Vulnerability Audit: SQL Injection & XSS vulnerability checks completed\n" \
                  f"  - SSL Certificate Validation: Completed"
        p4.font.name = 'Arial'
        p4.font.size = Pt(14)
        p4.font.color.rgb = text_dark
        p4.space_before = Pt(12)
        slide3 = prs.slides.add_slide(slide_layout)
        set_bg(slide3, light_bg)
        add_slide_header(slide3, "Vulnerability & Security Scan Results")
        results_box = slide3.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.833), Inches(5.0))
        tf3 = results_box.text_frame
        tf3.word_wrap = True
        tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
        p = tf3.paragraphs[0]
        p.text = "Key Scan Findings"
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = text_dark
        sql_vuln = any(item.get('vulnerable') for item in results.get('sql_injection', []))
        p2 = tf3.add_paragraph()
        p2.text = f"- SQL Injection Tester: {'POTENTIAL VULNERABILITY DETECTED' if sql_vuln else 'No Indicator Found'}"
        p2.font.name = 'Arial'
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(220, 38, 38) if sql_vuln else RGBColor(5, 150, 105)
        p2.space_before = Pt(12)
        xss_vuln = any(item.get('vulnerable') for item in results.get('xss', []))
        p3 = tf3.add_paragraph()
        p3.text = f"- XSS Detection: {'POTENTIAL VULNERABILITY DETECTED' if xss_vuln else 'No Indicator Found'}"
        p3.font.name = 'Arial'
        p3.font.size = Pt(15)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(220, 38, 38) if xss_vuln else RGBColor(5, 150, 105)
        p3.space_before = Pt(8)
        ssl_status = results.get('ssl', {})
        p4 = tf3.add_paragraph()
        p4.text = f"- SSL Checker: Issuer = {ssl_status.get('issuer', 'N/A')} | Valid = {ssl_status.get('valid', 'N/A')}"
        p4.font.name = 'Arial'
        p4.font.size = Pt(14)
        p4.font.color.rgb = text_dark
        p4.space_before = Pt(12)
        open_ports = [item for item in results.get('port_scan', []) if item.get('status') == 'Open']
        p5 = tf3.add_paragraph()
        p5.text = f"- Open Ports Detected: {', '.join(str(p['port']) for p in open_ports) if open_ports else 'None'}"
        p5.font.name = 'Arial'
        p5.font.size = Pt(14)
        p5.font.color.rgb = text_dark
        p5.space_before = Pt(8)
        slide4 = prs.slides.add_slide(slide_layout)
        set_bg(slide4, light_bg)
        add_slide_header(slide4, "Recommendations & Remediation")
        rec_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.833), Inches(5.0))
        tf4 = rec_box.text_frame
        tf4.word_wrap = True
        tf4.margin_left = tf4.margin_top = tf4.margin_right = tf4.margin_bottom = 0
        p = tf4.paragraphs[0]
        p.text = "Remediation Action Plan"
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = text_dark
        for item in results.get('recommendations', []):
            p_rec = tf4.add_paragraph()
            p_rec.text = f"- {item}"
            p_rec.font.name = 'Arial'
            p_rec.font.size = Pt(14)
            p_rec.font.color.rgb = text_dark
            p_rec.space_before = Pt(10)
        path = self.output_dir / f'{base_name}.pptx'
        prs.save(str(path))
        return path

    def save_pdf(self, results: dict[str, Any], base_name: str) -> Path | None:
        if SimpleDocTemplate is None:
            note_path = self.output_dir / f'{base_name}_pdf_not_created.txt'
            note_path.write_text('PDF report was not created because reportlab is not installed. Run: pip install -r requirements.txt', encoding='utf-8')
            return None
        path = self.output_dir / f'{base_name}.pdf'
        styles = getSampleStyleSheet()
        doc = SimpleDocTemplate(str(path), pagesize=A4)
        story: list[Any] = []
        story.append(Paragraph('Advanced VAPT Suite Security Scan Report', styles['Title']))
        story.append(Paragraph(f"Target: {results['target_url']}", styles['Normal']))
        story.append(Paragraph(f"Scan Time: {results['scan_time']}", styles['Normal']))
        story.append(Paragraph(f"Risk Score: {results['risk_score']}/100 ({results['risk_severity']})", styles['Heading2']))
        story.append(Spacer(1, 12))
        
        recon_data = results.get('recon_scan', {})
        story.append(Paragraph('Reconnaissance Scan', styles['Heading2']))
        story.append(Paragraph(f"<b>Fingerprinted Technologies:</b> {', '.join(recon_data.get('technologies', [])) or 'None'}", styles['Normal']))
        story.append(Spacer(1, 6))
        
        subdomains_list = [[item['subdomain'], item['ip']] for item in recon_data.get('subdomains', [])]
        if subdomains_list:
            story.extend(self._pdf_table('Discovered Subdomains', [['Subdomain', 'IP Address']] + subdomains_list[:12]))
        else:
            story.append(Paragraph('No active subdomains mapped.', styles['Normal']))
        story.append(Spacer(1, 12))

        story.extend(self._pdf_table('Port Scan', [['Port', 'Service', 'Status']] + [[row['port'], row['service'], row['status']] for row in results.get('port_scan', [])]))
        story.extend(self._pdf_table('SQL Injection', [['Parameter', 'Payload', 'Status']] + [[row['parameter'], row['payload'], 'Potential' if row['vulnerable'] else 'No indicator'] for row in results.get('sql_injection', [])[:8]]))
        story.extend(self._pdf_table('XSS', [['Parameter', 'Payload', 'Status']] + [[row['parameter'], row['payload'], 'Potential' if row['vulnerable'] else 'No indicator'] for row in results.get('xss', [])[:8]]))
        
        ssl_status = results.get('ssl', {})
        story.append(Paragraph('SSL Status', styles['Heading2']))
        story.append(Paragraph(f"Issuer: {ssl_status.get('issuer')}", styles['Normal']))
        story.append(Paragraph(f"Expiry: {ssl_status.get('expiry_date')}", styles['Normal']))
        
        story.append(Paragraph('Recommendations', styles['Heading2']))
        for item in results.get('recommendations', []):
            story.append(Paragraph(f'- {item}', styles['Normal']))
        doc.build(story)
        return path

    @staticmethod
    def _pdf_table(title: str, rows: list[list[Any]]) -> list[Any]:
        styles = getSampleStyleSheet()
        story: list[Any] = [Paragraph(title, styles['Heading2'])]
        table = Table(rows, repeatRows=1)
        table.setStyle(TableStyle([('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#102a43')), ('TEXTCOLOR', (0, 0), (-1, 0), colors.white), ('GRID', (0, 0), (-1, -1), 0.25, colors.grey), ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'), ('VALIGN', (0, 0), (-1, -1), 'TOP')]))
        story.append(table)
        story.append(Spacer(1, 12))
        return story